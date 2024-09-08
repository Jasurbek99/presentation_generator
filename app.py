import os
import random
import re
import time
import collections.abc as collections_abc
from dotenv import load_dotenv
from flask import Flask, render_template, request, redirect, url_for
from pptx import Presentation
from openai import OpenAI
import prompts
from pptx.enum.shapes import MSO_SHAPE_TYPE
app = Flask(__name__)

# Set your OpenAI API key
load_dotenv()
api_key = os.environ.get("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)
model_type = "gpt-4-turbo"  # Or "gpt-4" if you have access

def create_ppt_text(prompt, slides, info="",model_type='gpt-3.5-turbo'):
    final_prompt = prompts.make_prompt(prompt, slides, info, model_type)
    print(f"Model=={model_type}")
    # Call the OpenAI API
    try:
        response = client.chat.completions.create(
            model=model_type,
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": final_prompt}
            ],
            max_tokens=4096,
            temperature=0.6,
            top_p=0.95,
            n=1,
            stop=None
        )
        
        result = response.choices[0].message.content
        return "Title:" + result
    except Exception as e:
        print(f"Error in OpenAI API call: {e}")
        return "Title:Error in generating slide content"

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('Title:'):
                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):
                content = line.replace('Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
                
        if content:
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[slide_placeholder_index]
            tf = body_shape.text_frame
            tf.text = content
       
       
                            
    prs.save(f'static/GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"static/GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"

def choose_slide_layout(last_slide_layout_index, first_time):
    layout_indices = [1, 7, 8]
    if first_time:
        return 1
    slide_layout_index = last_slide_layout_index
    while slide_layout_index == last_slide_layout_index:
        slide_layout_index = random.choice(layout_indices)
    return slide_layout_index

def generate_ppt(prompt, add_info, slides, theme,model_type):
    prompt = re.sub(r'[^\w\s.\-\(\)]', '', prompt)
    if theme not in range(1, 8):
        print("Invalid theme number, default theme will be applied.")
        theme = 1
    
    print("Generating the PowerPoint, this could take some time depending on your GPU...\n")
    
    try:
        with open(f'Cache/{prompt}.txt', 'w', encoding='utf-8') as f:
            f.write(create_ppt_text(prompt, slides, add_info,model_type))
        ppt_path = create_ppt(f'Cache/{prompt}.txt', theme, prompt)
        return str(ppt_path)
    except IOError as e:
        print(f"Error creating PowerPoint file: {e}")
        return ""

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        topic = request.form['topic']
        add_info = request.form['add_info']
        slides = int(request.form['slides'])
        theme = int(request.form['theme'])
        language = request.form['language']
        print(language)
        print(f'theme={theme}')
        if language == 'Turkmen':
            model_type = "gpt-4o"
        else:
            model_type = "gpt-3.5-turbo"
        add_info=add_info+f" Presentation must be in {language} language "
        start_time = time.time()
        ppt_path = generate_ppt(topic, add_info, slides, theme,model_type)
        end_time = time.time()
       
        if ppt_path:
            elapsed_time = round((end_time - start_time), 2)
            return redirect(url_for('result', filepath=ppt_path, time=elapsed_time))
        else:
            message = "Failed to generate PowerPoint."
            return render_template('index.html', error=message)
    return render_template('index.html')

@app.route('/result')
def result():
    filepath = request.args['filepath']
    elapsed_time = request.args['time']
    return render_template('result.html', filepath=filepath, time=elapsed_time)

if __name__ == '__main__':
    app.run(debug=True)
