from flask import Flask, render_template, request, redirect, url_for
import os
import re
import random
import time
import requests
from pptx import Presentation
from pptx.util import Inches
import openai

app = Flask(__name__)
openai.api_key = os.getenv("OPENAI_API_KEY")

# DALL-E 2 API call
def generate_dalle_image(prompt):
    try:
        response = openai.images.generate(
            prompt=prompt,
            n=1,
            size="1024x1024"
        )
        image_url = response.data[0].url
        image_response = requests.get(image_url)
        if image_response.status_code == 200:
            return image_response.content
        else:
            return None
    except Exception as e:
        print(f"Error generating image: {e}")
        return None

def save_image(content, filename):
    with open(filename, 'wb') as f:
        f.write(content)

def create_ppt_text(prompt, slides, info=""):
    model_type = "gpt-4-turbo" if info else "gpt-3.5-turbo"
    final_prompt = f"{prompt} {slides} {info} {model_type}"
    
    try:
        response = openai.chat.completions.create(
            model=model_type,
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": final_prompt}
            ],
            max_tokens=4096,
            temperature=0.6,
            top_p=0.95,
            n=1,
            stop=None
        )
        
        result = response.choices[0].message.content
        return "Title:" + result
    except Exception as e:
        print(f"Error in OpenAI API call: {e}")
        return "Title:Error in generating slide content"

def create_ppt(text_file, design_number, ppt_name, image_filename=None):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('Title:'):
                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                
                if image_filename:
                    left = Inches(1)
                    top = Inches(1.5)
                    height = Inches(3.5)
                    slide.shapes.add_picture(image_filename, left, top, height=height)
                continue
            elif line.startswith('Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):
                content = line.replace('Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
                
        if content:
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[slide_placeholder_index]
            tf = body_shape.text_frame
            tf.text = content

    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"

def generate_ppt(prompt, add_info, slides, theme, model_type):
    prompt = re.sub(r'[^\w\s.\-\(\)]', '', prompt)
    if theme not in range(1, 8):
        print("Invalid theme number, default theme will be applied.")
        theme = 1

    print("Generating the PowerPoint, this could take some time depending on your GPU...\n")

    try:
        with open(f'Cache/{prompt}.txt', 'w', encoding='utf-8') as f:
            f.write(create_ppt_text(prompt, slides, add_info))
        
        # Generate DALL-E image for the title slide
        image_content = generate_dalle_image(prompt)
        if image_content:
            image_filename = f'Cache/{prompt}_image.png'
            save_image(image_content, image_filename)
        else:
            image_filename = None
        
        ppt_path = create_ppt(f'Cache/{prompt}.txt', theme, prompt, image_filename)
        return str(ppt_path)
    except IOError as e:
        print(f"Error creating PowerPoint file: {e}")
        return ""

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        language = request.form['language']
        topic = request.form['topic']
        add_info = request.form['add_info']
        slides = int(request.form['slides'])
        theme = int(request.form['theme'])
        
        # Adjust the model based on language
        if language == 'tk':
            model_type = "gpt-4-o"
        else:
            model_type = "gpt-3.5-turbo"

        # Generate the PowerPoint
        start_time = time.time()
        ppt_path = generate_ppt(topic, add_info, slides, theme, model_type)
        end_time = time.time()
        
        if ppt_path:
            elapsed_time = round((end_time - start_time), 2)
            return redirect(url_for('result', filepath=ppt_path, time=elapsed_time, language=language))
        else:
            message = "Failed to generate PowerPoint."
            return render_template('index.html', error=message)
    
    return render_template('index.html')

@app.route('/result')
def result():
    filepath = request.args['filepath']
    elapsed_time = request.args['time']
    language = request.args['language']
    relative_filepath = filepath.replace('static/', '')
    return render_template('result.html', filepath=relative_filepath, time=elapsed_time, language=language)

if __name__ == "__main__":
    app.run(debug=True)
